from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
TEMPLATE = Path("/Users/gh/Desktop/AI计算机/Generation-AI 中期汇报模版(6).pptx")
INPUT = ROOT / "midterm-draft.md"
OUTPUT = ROOT / "midterm.pptx"

FONT = "Times New Roman"
TITLE_SIZE = Pt(28)
BODY_SIZE = Pt(14)
REF_SIZE = Pt(12)
LINE_SPACING = 1.15

EXPECTED_SECTIONS = [
    "Title",
    "Research Question & Hypothesis",
    "Background",
    "Literature Review",
    "Research Design / Method",
    "Research Plan & Challenges",
    "Expected Results — user study not yet run",
    "References",
    "Acknowledgements",
]


def parse_sections(markdown: str) -> dict[str, str]:
    matches = list(re.finditer(r"^##\s+(.+)$", markdown, flags=re.MULTILINE))
    sections = {}
    for index, match in enumerate(matches):
        title = match.group(1).strip()
        start = match.end()
        end = matches[index + 1].start() if index + 1 < len(matches) else len(markdown)
        sections[title] = markdown[start:end].strip()

    missing = [section for section in EXPECTED_SECTIONS if not sections.get(section)]
    if missing:
        raise ValueError("Missing or empty section(s): " + ", ".join(missing))
    return sections


def split_paragraphs(text: str) -> list[str]:
    return [p.strip().replace("\n", " ") for p in re.split(r"\n\s*\n", text) if p.strip()]


def lines(text: str) -> list[str]:
    return [line.strip() for line in text.splitlines() if line.strip()]


def delete_slide(prs: Presentation, slide_index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[slide_index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    slide_id_list.remove(slide_id)


def clear_text_shapes(slide) -> None:
    for shape in list(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            shape.text_frame.clear()
            shape.element.getparent().remove(shape.element)


def set_paragraph_style(paragraph, font_size=BODY_SIZE, bold=False, color=RGBColor(28, 36, 48)):
    paragraph.line_spacing = LINE_SPACING
    for run in paragraph.runs:
        run.font.name = FONT
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color


def add_textbox(slide, text, x, y, w, h, font_size=BODY_SIZE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)

    for idx, part in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = part
        paragraph.alignment = align
        set_paragraph_style(paragraph, font_size=font_size, bold=bold)
    return box


def add_title(slide, title):
    add_textbox(slide, title, 0.65, 0.35, 12.0, 0.55, font_size=TITLE_SIZE, bold=True)


def bullet_text(items: list[str]) -> str:
    return "\n".join(f"• {item}" for item in items)


def first_sentence(text: str) -> str:
    match = re.match(r"(.+?[.!?])\s", text)
    return match.group(1) if match else text


def write_title_slide(slide, content: str) -> None:
    clear_text_shapes(slide)
    title_lines = lines(content)
    add_textbox(
        slide,
        title_lines[0],
        1.2,
        2.15,
        10.9,
        0.9,
        font_size=TITLE_SIZE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        "\n".join(title_lines[1:3]),
        2.0,
        3.35,
        9.3,
        0.9,
        font_size=BODY_SIZE,
        align=PP_ALIGN.CENTER,
    )


def write_standard_slide(slide, title: str, body: str) -> None:
    clear_text_shapes(slide)
    add_title(slide, title)
    add_textbox(slide, body, 0.8, 1.15, 11.7, 5.65, font_size=BODY_SIZE)


def write_expected_slide(slide, body_items: list[str]) -> None:
    clear_text_shapes(slide)
    add_title(slide, "Expected Results — user study not yet run")
    add_textbox(slide, bullet_text(body_items), 0.7, 1.15, 5.8, 5.7, font_size=BODY_SIZE)
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.0),
        Inches(1.25),
        Inches(5.55),
        Inches(4.95),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(243, 244, 246)
    rect.line.color.rgb = RGBColor(145, 152, 161)
    rect.line.width = Pt(1.5)
    rect.line.dash_style = 4
    add_textbox(
        slide,
        "Insert prototype screenshot here",
        7.35,
        3.35,
        4.85,
        0.5,
        font_size=BODY_SIZE,
        align=PP_ALIGN.CENTER,
    )


def write_references_slide(slide, refs: list[str]) -> None:
    clear_text_shapes(slide)
    add_title(slide, "References")
    split = (len(refs) + 1) // 2
    add_textbox(slide, "\n".join(refs[:split]), 0.55, 1.05, 6.05, 5.95, font_size=REF_SIZE)
    add_textbox(slide, "\n".join(refs[split:]), 6.85, 1.05, 5.95, 5.95, font_size=REF_SIZE)


def main() -> None:
    sections = parse_sections(INPUT.read_text(encoding="utf-8"))
    prs = Presentation(str(TEMPLATE))

    # The template's first slide is an instruction page, not part of the final deck.
    delete_slide(prs, 0)
    if len(prs.slides) != 9:
        raise ValueError(f"Expected 9 content slides after removing instructions, got {len(prs.slides)}")

    write_title_slide(prs.slides[0], sections["Title"])
    write_standard_slide(prs.slides[1], "Research Question & Hypothesis", sections["Research Question & Hypothesis"])

    background_items = [first_sentence(p) for p in split_paragraphs(sections["Background"])]
    write_standard_slide(prs.slides[2], "Background", bullet_text(background_items))

    literature_items = [
        "Zhao et al. (2024): clean heating costs are highly sensitive to subsidies.",
        "Yu and Xin (2021): heating burden can become a form of energy poverty.",
        "Li, Song, and Zhu (2021): gas costs can change household heating choices.",
        "He et al. (2021): some Baoding clean-heating households returned to coal.",
        "Zhang et al. (2024): costs and health benefits are uneven across regions.",
        "My gap: a village-facing sandbox based on household data.",
    ]
    write_standard_slide(prs.slides[3], "Literature Review", bullet_text(literature_items))

    method_items = [
        "The artifact is a clean heating transition simulation sandbox.",
        "Users enter income, house size, heating method, and winter heating cost.",
        "The current MVP shows cost, surplus, compliance, emissions, and energy burden.",
        "The next version will use many households from one village.",
        "Feedback will come from short before-and-after questions and user comments.",
    ]
    write_standard_slide(prs.slides[4], "Research Design / Method", bullet_text(method_items))

    plan_items = [
        "June 20: present the research question, hypothesis, MVP, and evidence.",
        "Late June: improve simulator logic and pathway recommendations.",
        "Early July: build a simple village data table.",
        "After that: run a small user test and collect feedback.",
        "Main challenges: data privacy, model accuracy, and clear communication.",
    ]
    write_standard_slide(prs.slides[5], "Research Plan & Challenges", bullet_text(plan_items))

    expected_items = [
        "I expect users to see that there is no single best pathway for every village.",
        "I expect affordability to become more visible.",
        "I expect users to care more about household-data-based village decisions.",
        "I will not report user numbers or final conclusions before the study is run.",
    ]
    write_expected_slide(prs.slides[6], expected_items)

    write_references_slide(prs.slides[7], lines(sections["References"]))
    write_standard_slide(prs.slides[8], "Acknowledgements", sections["Acknowledgements"])

    prs.save(str(OUTPUT))
    print(f"Generated {OUTPUT}")


if __name__ == "__main__":
    main()
